"""Structural QA on the generated .pptx -- catches issues without rendering."""
import sys
from pptx import Presentation
from pptx.util import Emu

p = Presentation("D:/Claude/auto_cst/runs/auto_cst_NIR_presentation.pptx")
sw = Emu(p.slide_width).inches
sh = Emu(p.slide_height).inches

print(f"Slide canvas: {sw:.2f}\" x {sh:.2f}\"")
print(f"Slide count: {len(p.slides)}\n")

issues_total = 0

def emu_to_in(e): return Emu(e).inches if e is not None else None

for i, slide in enumerate(p.slides, 1):
    print(f"=== Slide {i} ===")
    boxes = []
    for s in slide.shapes:
        try:
            x = emu_to_in(s.left)
            y = emu_to_in(s.top)
            w = emu_to_in(s.width)
            h = emu_to_in(s.height)
        except Exception:
            continue
        if x is None or y is None or w is None or h is None:
            continue
        text = ""
        if s.has_text_frame:
            text = s.text_frame.text.strip()
        boxes.append((s.shape_id, x, y, w, h, text[:60]))

    issues = []

    # 1. Off-canvas check
    for sid, x, y, w, h, t in boxes:
        if x < 0 or y < 0:
            issues.append(f"  [OFF-CANVAS]  shape {sid} starts at ({x:.2f},{y:.2f}): {t!r}")
        if x + w > sw + 0.01:
            issues.append(f"  [OFF-CANVAS]  shape {sid} extends to x={x+w:.2f}\" (canvas {sw:.2f}\"): {t!r}")
        if y + h > sh + 0.01:
            issues.append(f"  [OFF-CANVAS]  shape {sid} extends to y={y+h:.2f}\" (canvas {sh:.2f}\"): {t!r}")

    # 2. Margin check (>= 0.5" recommended for most content; allow accent bars at x=0)
    for sid, x, y, w, h, t in boxes:
        # Skip accent bars (very thin shapes at x=0)
        if x == 0 and w < 0.2:
            continue
        # Skip background-coloured rectangles that span full width (header/footer bars)
        if x < 0.1 and w > sw - 0.2:
            continue
        if t and (x < 0.4 or y < 0.2):
            issues.append(f"  [TIGHT MARGIN]  shape {sid} at ({x:.2f},{y:.2f}): {t!r}")
        if t and (x + w > sw - 0.4 or y + h > sh - 0.2):
            # only flag if content has text; bars/accents OK
            issues.append(f"  [TIGHT MARGIN]  shape {sid} ends at ({x+w:.2f},{y+h:.2f}): {t!r}")

    # 3. Pairwise overlap check among text-bearing shapes (excluding background fills)
    text_boxes = [b for b in boxes if b[5]]
    for j in range(len(text_boxes)):
        for k in range(j + 1, len(text_boxes)):
            a = text_boxes[j]; b = text_boxes[k]
            if (a[1] < b[1] + b[3] and a[1] + a[3] > b[1] and
                a[2] < b[2] + b[4] and a[2] + a[4] > b[2]):
                # Compute overlap area
                ox = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                oy = min(a[2] + a[4], b[2] + b[4]) - max(a[2], b[2])
                area = ox * oy
                if area > 0.05:   # small overlap is fine (text in a card with bg)
                    issues.append(f"  [OVERLAP]    {a[0]} <{a[5]!r:.40}> overlaps {b[0]} <{b[5]!r:.40}> by {area:.2f} sq.in")

    if issues:
        for issue in issues[:8]:  # cap output per slide
            print(issue)
        if len(issues) > 8:
            print(f"  ... and {len(issues) - 8} more")
        issues_total += len(issues)
    else:
        print("  OK")
    print()

print(f"\nTotal issues: {issues_total}")
sys.exit(0 if issues_total == 0 else 1)
